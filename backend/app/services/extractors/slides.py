"""
SlidesExtractor — PF-03 Interactive Slides.

Converts .pptx to a series of slide images and extracts text for AI processing.

Flow:
  1. Download / read .pptx bytes from source_url (file:// local upload)
  2. Use python-pptx to extract text from each slide (for AI outputs + RAG)
  3. Use LibreOffice headless to convert .pptx → .pdf
  4. Use PyMuPDF (fitz) to render each PDF page to a PNG image
  5. Store images in /data/slide_outputs/{content_item_id}/slide_{n}.png
  6. Return ExtractedContent with full text + slide_assets metadata stored separately

The pipeline will call update_slide_assets() to persist image paths to
content_items.slide_assets after extraction.

LibreOffice must be installed:
    sudo snap install libreoffice
"""
from __future__ import annotations

import os
import subprocess
import tempfile
import uuid
from pathlib import Path

import fitz  # PyMuPDF
import structlog
from pptx import Presentation
from pptx.util import Pt

from app.core.exceptions import ContentProcessingError
from app.utils.hashing import sha256_bytes, sha256_text
from .base import BaseExtractor, ExtractedContent

log = structlog.get_logger(__name__)

SLIDE_OUTPUT_DIR = Path(os.environ.get("SLIDE_OUTPUT_DIR", "/data/slide_outputs"))
LIBREOFFICE_CMD = os.environ.get("LIBREOFFICE_CMD", "libreoffice")

# Max PPTX size: 200MB
MAX_PPTX_BYTES = 200 * 1024 * 1024
# Max image width for slide thumbnails (px)
SLIDE_IMAGE_WIDTH = 1280


class SlidesExtractor(BaseExtractor):
    """Extracts text and renders slide images from .pptx files."""

    @property
    def supported_content_types(self) -> list[str]:
        return ["interactive_slides"]

    async def extract(
        self,
        *,
        url: str | None = None,
        file_bytes: bytes | None = None,
        content_item_metadata: dict | None = None,
    ) -> ExtractedContent:
        meta = content_item_metadata or {}
        content_item_id = meta.get("content_item_id") or str(uuid.uuid4())

        # Acquire PPTX bytes
        if file_bytes is None:
            if url is None:
                raise ContentProcessingError("SlidesExtractor requires url or file_bytes")
            file_bytes = await self._download_file(url)

        if len(file_bytes) > MAX_PPTX_BYTES:
            raise ContentProcessingError(f"PPTX too large: {len(file_bytes) // 1024 // 1024}MB")

        content_hash = sha256_bytes(file_bytes)

        # ── Extract text from slides ──────────────────────────────────────
        slide_texts, slide_count = _extract_slide_texts(file_bytes)
        full_text = _join_slide_texts(slide_texts)
        word_count = len(full_text.split())

        log.info("slides_text_extracted", slide_count=slide_count, word_count=word_count)

        # ── Convert to images ─────────────────────────────────────────────
        out_dir = SLIDE_OUTPUT_DIR / content_item_id
        out_dir.mkdir(parents=True, exist_ok=True)

        slide_assets = _convert_pptx_to_images(file_bytes, out_dir, slide_count)

        log.info("slides_images_rendered", slide_count=len(slide_assets), dir=str(out_dir))

        return ExtractedContent(
            raw_text=full_text,
            content_hash=content_hash,
            page_count=slide_count,
            word_count=word_count,
            segments=[],
            all_segments={},
            detected_source_language=None,
            extraction_metadata={
                "slide_count": slide_count,
                "slide_texts": slide_texts,
                "slide_assets": slide_assets,  # pipeline will persist this
                "extractor": "pptx+libreoffice",
            },
        )

    async def _download_file(self, url: str) -> bytes:
        """Download from file:// or http(s)://."""
        if url.startswith("file://"):
            path = Path(url[7:])
            if not path.exists():
                raise ContentProcessingError(f"File not found: {path}")
            return path.read_bytes()

        import httpx
        timeout = httpx.Timeout(connect=10.0, read=120.0, write=10.0, pool=5.0)
        async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
            resp = await client.get(url)
            resp.raise_for_status()
            return resp.content


def _extract_slide_texts(pptx_bytes: bytes) -> tuple[list[str], int]:
    """Extract text from each slide using python-pptx."""
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(pptx_bytes)
        tmp_path = Path(tmp.name)

    try:
        prs = Presentation(tmp_path)
        slide_texts: list[str] = []

        for slide_idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text.strip())
                    if line.strip():
                        texts.append(line.strip())
            slide_texts.append("\n".join(texts) if texts else f"[Slide {slide_idx + 1} — no text]")

        return slide_texts, len(prs.slides)
    finally:
        tmp_path.unlink(missing_ok=True)


def _join_slide_texts(slide_texts: list[str]) -> str:
    """Format all slide texts as a single document for AI processing."""
    parts = []
    for i, text in enumerate(slide_texts, 1):
        parts.append(f"## Slide {i}\n\n{text}")
    return "\n\n".join(parts)


def _convert_pptx_to_images(
    pptx_bytes: bytes,
    out_dir: Path,
    slide_count: int,
) -> list[dict]:
    """
    Convert PPTX to slide images using LibreOffice → PDF → PyMuPDF.

    Returns list of {index, path, thumbnail_path, width, height}.
    """
    slide_assets: list[dict] = []

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_pptx = Path(tmpdir) / "presentation.pptx"
        tmp_pptx.write_bytes(pptx_bytes)

        # LibreOffice headless convert to PDF
        tmp_pdf = Path(tmpdir) / "presentation.pdf"
        try:
            result = subprocess.run(
                [
                    LIBREOFFICE_CMD,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", tmpdir,
                    str(tmp_pptx),
                ],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                log.warning(
                    "libreoffice_conversion_failed",
                    stdout=result.stdout[:500],
                    stderr=result.stderr[:500],
                )
                # Fallback: create placeholder images
                return _create_placeholder_images(slide_count, out_dir)
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            log.warning("libreoffice_unavailable", error=str(e))
            return _create_placeholder_images(slide_count, out_dir)

        if not tmp_pdf.exists():
            log.warning("pdf_not_generated")
            return _create_placeholder_images(slide_count, out_dir)

        # Render PDF pages → PNG using PyMuPDF
        doc = fitz.open(str(tmp_pdf))
        actual_count = doc.page_count

        for page_num in range(actual_count):
            page = doc[page_num]
            # Scale to SLIDE_IMAGE_WIDTH px wide
            scale = SLIDE_IMAGE_WIDTH / page.rect.width
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            slide_index = page_num + 1  # 1-based
            image_path = out_dir / f"slide_{slide_index:03d}.png"
            thumb_path = out_dir / f"slide_{slide_index:03d}_thumb.jpg"

            pix.save(str(image_path))

            # Thumbnail at 320px wide
            thumb_scale = 320 / page.rect.width
            thumb_mat = fitz.Matrix(thumb_scale, thumb_scale)
            thumb_pix = page.get_pixmap(matrix=thumb_mat, alpha=False)
            thumb_pix.save(str(thumb_path))

            slide_assets.append({
                "index": slide_index,
                "path": str(image_path),
                "thumbnail_path": str(thumb_path),
                "width": pix.width,
                "height": pix.height,
            })

        doc.close()

    return slide_assets


def _create_placeholder_images(slide_count: int, out_dir: Path) -> list[dict]:
    """
    When LibreOffice is unavailable, generate grey placeholder images
    so the slide player still works (will show slide number + text overlay).
    """
    assets = []
    for i in range(1, slide_count + 1):
        # Create a simple 1280×720 grey image with slide number
        page = fitz.open()
        page_obj = page.new_page(width=1280, height=720)
        page_obj.draw_rect(fitz.Rect(0, 0, 1280, 720), fill=(0.9, 0.9, 0.9))
        page_obj.insert_text(
            (560, 360),
            f"Slide {i}",
            fontsize=48,
            color=(0.4, 0.4, 0.4),
        )

        pix = page_obj.get_pixmap()
        image_path = out_dir / f"slide_{i:03d}.png"
        pix.save(str(image_path))

        thumb_path = out_dir / f"slide_{i:03d}_thumb.jpg"
        # Create small thumb
        mat = fitz.Matrix(0.25, 0.25)
        thumb_pix = page_obj.get_pixmap(matrix=mat)
        thumb_pix.save(str(thumb_path))

        assets.append({
            "index": i,
            "path": str(image_path),
            "thumbnail_path": str(thumb_path),
            "width": 1280,
            "height": 720,
        })
        page.close()

    return assets
